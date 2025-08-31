"""
Flexible content processor for various file types
Extensible and minimal intrusion design
"""

from typing import Dict, Any, List, Optional, Tuple
import mimetypes
import os
from pathlib import Path


class ContentProcessor:
    """
    Process various content types into indexable format
    Designed for flexibility and extensibility
    """
    
    def __init__(self):
        # Register processors by mime type patterns
        self.processors = {}
        self._register_default_processors()
    
    def _register_default_processors(self):
        """Register built-in processors"""
        # Text-based files
        self.register_processor(
            patterns=["text/*", "application/json", "application/xml"],
            processor=self._process_text
        )
        
        # Code files (often detected as text/plain)
        self.register_processor(
            extensions=[".py", ".js", ".java", ".cpp", ".c", ".rs", ".go"],
            processor=self._process_code
        )
        
        # Markdown
        self.register_processor(
            extensions=[".md", ".markdown"],
            processor=self._process_markdown
        )
        
        # PDF (if libraries available)
        self.register_processor(
            patterns=["application/pdf"],
            processor=self._process_pdf
        )
        
        # Microsoft Office documents
        self.register_processor(
            extensions=[".docx", ".doc"],
            processor=self._process_docx
        )
        self.register_processor(
            extensions=[".xlsx", ".xls"],
            processor=self._process_excel
        )
        self.register_processor(
            extensions=[".pptx", ".ppt"],
            processor=self._process_powerpoint
        )
        
        # Images (if libraries available)
        self.register_processor(
            patterns=["image/*"],
            processor=self._process_image
        )
    
    def register_processor(
        self,
        patterns: List[str] = None,
        extensions: List[str] = None,
        processor: callable = None
    ):
        """
        Register a content processor
        Allows external registration for extensibility
        """
        if patterns:
            for pattern in patterns:
                self.processors[pattern] = processor
        if extensions:
            for ext in extensions:
                self.processors[f"ext:{ext}"] = processor
    
    async def process_content(
        self,
        source: str,
        source_type: str = "file",
        **kwargs
    ) -> Tuple[str, Dict[str, Any]]:
        """
        Process content from various sources
        Returns: (processed_text, metadata)
        """
        if source_type == "text":
            # Direct text input
            return source, {"type": "text", "source": "direct_input"}
        
        elif source_type == "file":
            # Detect file type and process
            return await self._process_file(source, **kwargs)
        
        elif source_type == "url":
            # Web content (future extension)
            return await self._process_url(source, **kwargs)
        
        else:
            # Unknown type, return as-is
            return str(source), {"type": "unknown"}
    
    async def _process_file(self, file_path: str, **kwargs) -> Tuple[str, Dict[str, Any]]:
        """Process file based on type detection"""
        path = Path(file_path)
        
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        # Detect mime type
        mime_type, _ = mimetypes.guess_type(file_path)
        extension = path.suffix.lower()
        
        # Build metadata
        metadata = {
            "source": file_path,
            "mime_type": mime_type,
            "extension": extension,
            "size": path.stat().st_size,
            "modified": path.stat().st_mtime
        }
        
        # Find appropriate processor
        processor = self._find_processor(mime_type, extension)
        
        # Process content
        try:
            content = await processor(file_path, metadata, **kwargs)
            return content, metadata
        except Exception as e:
            # Fallback to basic text reading
            print(f"Processor failed for {file_path}: {e}, using fallback")
            return await self._fallback_processor(file_path, metadata)
    
    def _find_processor(self, mime_type: str, extension: str):
        """Find appropriate processor for file type"""
        # Check extension first (more specific)
        ext_key = f"ext:{extension}"
        if ext_key in self.processors:
            return self.processors[ext_key]
        
        # Check exact mime type
        if mime_type and mime_type in self.processors:
            return self.processors[mime_type]
        
        # Check mime type patterns
        if mime_type:
            for pattern, processor in self.processors.items():
                if pattern.endswith("/*") and mime_type.startswith(pattern[:-2]):
                    return processor
        
        # Default processor
        return self._fallback_processor
    
    async def _process_text(self, file_path: str, metadata: Dict, **kwargs) -> str:
        """Process plain text files"""
        encoding = kwargs.get("encoding", "utf-8")
        try:
            with open(file_path, "r", encoding=encoding) as f:
                content = f.read()
            metadata["encoding"] = encoding
            return content
        except UnicodeDecodeError:
            # Try with different encoding
            with open(file_path, "r", encoding="latin-1") as f:
                content = f.read()
            metadata["encoding"] = "latin-1"
            return content
    
    async def _process_markdown(self, file_path: str, metadata: Dict, **kwargs) -> str:
        """Process markdown files"""
        content = await self._process_text(file_path, metadata, **kwargs)
        metadata["type"] = "markdown"
        
        # Optional: extract structure
        if kwargs.get("extract_structure", False):
            import re
            headers = re.findall(r'^#{1,6}\s+(.+)$', content, re.MULTILINE)
            metadata["headers"] = headers
        
        return content
    
    async def _process_code(self, file_path: str, metadata: Dict, **kwargs) -> str:
        """Process source code files"""
        content = await self._process_text(file_path, metadata, **kwargs)
        metadata["type"] = "code"
        
        # Optional: extract structure
        if kwargs.get("extract_structure", False):
            import re
            # Extract function/class definitions (basic)
            if metadata["extension"] in [".py"]:
                functions = re.findall(r'^def\s+(\w+)', content, re.MULTILINE)
                classes = re.findall(r'^class\s+(\w+)', content, re.MULTILINE)
                metadata["functions"] = functions
                metadata["classes"] = classes
        
        return content
    
    async def _process_pdf(self, file_path: str, metadata: Dict, **kwargs) -> str:
        """Process PDF files if library available"""
        try:
            import PyPDF2
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
            metadata["type"] = "pdf"
            metadata["pages"] = len(reader.pages)
            return text
        except ImportError:
            # Library not available, try basic extraction
            print("PyPDF2 not installed, using basic text extraction")
            return await self._fallback_processor(file_path, metadata)
    
    async def _process_docx(self, file_path: str, metadata: Dict, **kwargs) -> str:
        """Process Word DOCX files"""
        metadata["type"] = "docx"
        
        try:
            # Try python-docx first
            from docx import Document
            doc = Document(file_path)
            
            # Extract all paragraphs
            text = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text.append(paragraph.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text.append(cell.text)
            
            metadata["paragraphs"] = len(doc.paragraphs)
            metadata["tables"] = len(doc.tables)
            return "\n".join(text)
            
        except ImportError:
            # Try alternative method with zipfile
            try:
                import zipfile
                import xml.etree.ElementTree as ET
                
                text = []
                with zipfile.ZipFile(file_path, 'r') as z:
                    # Read main document
                    with z.open('word/document.xml') as f:
                        tree = ET.parse(f)
                        root = tree.getroot()
                        
                        # Extract text from XML
                        namespace = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
                        for elem in root.iter(namespace + 't'):
                            if elem.text:
                                text.append(elem.text)
                
                return " ".join(text)
                
            except Exception as e:
                print(f"Failed to extract DOCX content: {e}")
                return await self._fallback_processor(file_path, metadata)
    
    async def _process_excel(self, file_path: str, metadata: Dict, **kwargs) -> str:
        """Process Excel files"""
        metadata["type"] = "excel"
        
        try:
            import pandas as pd
            # Read all sheets
            dfs = pd.read_excel(file_path, sheet_name=None)
            
            text = []
            for sheet_name, df in dfs.items():
                text.append(f"Sheet: {sheet_name}")
                text.append(df.to_string())
            
            metadata["sheets"] = len(dfs)
            return "\n\n".join(text)
            
        except ImportError:
            print("pandas not installed for Excel processing")
            return await self._fallback_processor(file_path, metadata)
    
    async def _process_powerpoint(self, file_path: str, metadata: Dict, **kwargs) -> str:
        """Process PowerPoint files"""
        metadata["type"] = "powerpoint"
        
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            
            text = []
            for i, slide in enumerate(prs.slides, 1):
                text.append(f"Slide {i}:")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text.append(shape.text)
            
            metadata["slides"] = len(prs.slides)
            return "\n".join(text)
            
        except ImportError:
            print("python-pptx not installed for PowerPoint processing")
            return await self._fallback_processor(file_path, metadata)
    
    async def _process_image(self, file_path: str, metadata: Dict, **kwargs) -> str:
        """Process image files if OCR available"""
        metadata["type"] = "image"
        
        # Option 1: OCR extraction
        if kwargs.get("use_ocr", False):
            try:
                import pytesseract
                from PIL import Image
                img = Image.open(file_path)
                text = pytesseract.image_to_string(img)
                metadata["ocr"] = True
                return text
            except ImportError:
                pass
        
        # Option 2: Return description placeholder
        # In real implementation, could use vision model
        description = f"Image file: {Path(file_path).name}"
        
        # Extract basic image info if PIL available
        try:
            from PIL import Image
            img = Image.open(file_path)
            metadata["dimensions"] = img.size
            metadata["mode"] = img.mode
            description += f" ({img.size[0]}x{img.size[1]} {img.mode})"
        except ImportError:
            pass
        
        return description
    
    async def _process_url(self, url: str, metadata: Dict, **kwargs) -> str:
        """Process web content (future extension)"""
        metadata["type"] = "url"
        metadata["source"] = url
        # Implementation would fetch and parse web content
        return f"URL content: {url}"
    
    async def _fallback_processor(self, file_path: str, metadata: Dict) -> str:
        """Fallback for unknown file types"""
        # Try to read as text
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
            metadata["processor"] = "fallback_text"
            return content
        except:
            # Return file info as content
            metadata["processor"] = "fallback_info"
            return f"Binary file: {Path(file_path).name} ({metadata.get('size', 0)} bytes)"


# Singleton instance
_content_processor: Optional[ContentProcessor] = None


def get_content_processor() -> ContentProcessor:
    """Get or create content processor singleton"""
    global _content_processor
    if _content_processor is None:
        _content_processor = ContentProcessor()
    return _content_processor